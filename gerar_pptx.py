import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Tensor Mage Arena"
subtitle.text = "Aprendizado por Reforço e Redes Neurais no Jogo-MAGO-RL"

# Slide 2: O que é o Projeto?
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "O que é o Projeto?"
tf = content.text_frame
tf.text = "Ambiente de Aprendizado por Reforço (RL) para treinar um agente."
p = tf.add_paragraph()
p.text = "Agente (Mago) aprende através de Self-Play contínuo."
p = tf.add_paragraph()
p.text = "Objetivo: Ensinar heurísticas avançadas de combate (ex: esquivar, perseguir)."
p = tf.add_paragraph()
p.text = "Construído em JAX e Flax para processamento paralelo ultra-rápido."

# Slide 3: O que são Redes Neurais?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "O que são Redes Neurais?"
tf = content.text_frame
tf.text = "Modelos computacionais inspirados no cérebro biológico."
p = tf.add_paragraph()
p.text = "Compostas por camadas de neurônios artificiais interconectados."
p = tf.add_paragraph()
p.text = "Aprendem padrões em dados ajustando os pesos das conexões (treinamento)."
p = tf.add_paragraph()
p.text = "Permitem que a máquina 'veja' e tome decisões baseadas em complexidade matemática."

# Slide 4: Redes Neurais no Projeto
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Redes Neurais no Tensor Mage Arena"
tf = content.text_frame
tf.text = "Backbone Convolucional (CNN / ResNet): Processa a visão do jogo (mapa, inimigos)."
p = tf.add_paragraph()
p.text = "Conexões Residuais previnem perda de informação em redes profundas."
p = tf.add_paragraph()
p.text = "Memória Episódica (Self-Attention / Transformers): Processa 8 frames no tempo."
p = tf.add_paragraph()
p.text = "Em vez de LSTMs, o Transformer cruza passado e presente de imediato para prever física (tiros)."

# Slide 5: Como o Agente Aprende? (Reinforcement Learning)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Como o Agente Aprende? (Reinforcement Learning)"
tf = content.text_frame
tf.text = "Reward Shaping: O agente ganha pontos ao acertar e perde ao errar."
p = tf.add_paragraph()
p.text = "Recompensas densas por frame, não apenas ao vencer/perder a partida."
p = tf.add_paragraph()
p.text = "Punição de Tempo: Incentiva o agente a vencer rápido."
p = tf.add_paragraph()
p.text = "Punição de Campismo: Ficar parado é punido para forçar a exploração."
p = tf.add_paragraph()
p.text = "Punição de Ameaça: Ficar na reta de um tiro penaliza severamente (ensina o Dodging)."

# Slide 6: A Engine de Simulação
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "A Engine de Simulação"
tf = content.text_frame
tf.text = "Desenvolvida puramente com matrizes (Numpy/JAX) sem interfaces gráficas pesadas (Headless)."
p = tf.add_paragraph()
p.text = "Observação Espacial: Matriz que mistura paredes, jogadores e projéteis ativos."
p = tf.add_paragraph()
p.text = "Observação Escalar: Sensores de HP, Cooldowns e direção sugerida."
p = tf.add_paragraph()
p.text = "Ações Multimodais: Anda em 4 direções e mira num espaço 360 graus."

# Slide 7: Curriculum Learning
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Progressão Suave (Curriculum Learning)"
tf = content.text_frame
tf.text = "O aprendizado não tem fases fixas, possui uma dificuldade escalar que aumenta aos poucos."
p = tf.add_paragraph()
p.text = "A saúde, dano e velocidade inimiga sofrem 'Interpolação Linear' com a evolução do treino."
p = tf.add_paragraph()
p.text = "O Inimigo usa o Algoritmo A* com mapa de ameaças, fugindo dos tiros do agente."
p = tf.add_paragraph()
p.text = "O agente treina contra versões passadas de si mesmo (Self-Play Assíncrono)."

# Slide 8: Algoritmo PPO (Proximal Policy Optimization)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "O Algoritmo: PPO"
tf = content.text_frame
tf.text = "Mesmo algoritmo base por trás do ChatGPT da OpenAI."
p = tf.add_paragraph()
p.text = "Arquitetura Actor-Critic: O Actor decide as ações, o Critic julga quão bom o estado é."
p = tf.add_paragraph()
p.text = "Função de Custo Clipada: Impede que a rede sofra atualizações bruscas que a façam 'esquecer' o que aprendeu."

# Slide 9: Conclusão
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusão"
tf = content.text_frame
tf.text = "Tensor Mage Arena é um laboratório complexo de IA."
p = tf.add_paragraph()
p.text = "Uso avançado de Transformers e processamento matricial massivo (JAX)."
p = tf.add_paragraph()
p.text = "Táticas avançadas emergem naturalmente das leis da física e da pontuação imposta no treinamento."

prs.save('/home/jeanlaurentstoco/Desktop/prog1/jogo/Apresentacao_Tensor_Mage_Arena.pptx')
print("Apresentação gerada com sucesso!")
